from pptx import Presentation

SLD_LAYOUT_TITLE_AND_CONTENT = 1

prs = Presentation()
slide_layout = prs.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
slide = prs.slides.add_slide(slide_layout)

prs.save('my_first_ppt.pptx')